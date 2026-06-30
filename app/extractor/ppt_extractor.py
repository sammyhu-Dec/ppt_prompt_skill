from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas.models import ExtractedDeck, SlideExtracted


def _collect_shape_text(shape) -> List[str]:
    """递归提取 shape 文本。第一版重点处理文本框和组合形状。"""
    texts = []

    if hasattr(shape, "text") and shape.text:
        value = shape.text.strip()
        if value:
            texts.append(value)

    # 组合形状：递归读取子元素
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            texts.extend(_collect_shape_text(sub_shape))

    return texts


def _count_images(slide) -> int:
    count = 0
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            count += 1
    return count


def extract_ppt_text(pptx_path: Path) -> ExtractedDeck:
    """
    第一版 PPT 提取器：
    - 支持 .pptx
    - 提取每页文本
    - 粗略提取标题：取第一页非空文本
    - 统计图片数量

    暂不处理：
    - OCR
    - 图表深度解析
    - 页面截图理解
    - speaker notes
    """
    prs = Presentation(str(pptx_path))
    slides = []

    for page_idx, slide in enumerate(prs.slides, start=1):
        text_blocks: List[str] = []

        for shape in slide.shapes:
            text_blocks.extend(_collect_shape_text(shape))

        # 去重但保留顺序
        deduped = []
        seen = set()
        for text in text_blocks:
            if text not in seen:
                deduped.append(text)
                seen.add(text)

        raw_text = "\n".join(deduped).strip()
        title = deduped[0] if deduped else ""

        slides.append(
            SlideExtracted(
                page=page_idx,
                title=title,
                raw_text=raw_text,
                image_count=_count_images(slide),
                notes="",
            )
        )

    return ExtractedDeck(
        file_name=pptx_path.name,
        slide_count=len(slides),
        slides=slides,
    )
